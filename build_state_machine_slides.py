"""Generate state_machine_controller_slides.pptx for the robotics presentation.

Run: ../.venv/bin/python build_state_machine_slides.py
Source content mirrors state_machine_controller_slides.md.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ---------------------------------------------------------------
DARK = RGBColor(0x1F, 0x2A, 0x44)      # title / text
GREY = RGBColor(0x5B, 0x63, 0x72)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2E, 0x6F, 0xD6)      # joint
GREEN = RGBColor(0x2E, 0xA0, 0x5A)     # operational
RED = RGBColor(0xD6, 0x3B, 0x3B)       # force
ACCENT = RGBColor(0x6B, 0x4E, 0xC4)    # machine / misc

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill if fill else WHITE
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4):
    """runs: list of paragraphs; each para is list of (text, size, color, bold)."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (text, size, color, bold) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def bg(s, color=WHITE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def header(s, title, kicker=None):
    bar = box(s, 0, 0, SW, Inches(1.15), fill=DARK)
    txt(s, Inches(0.55), Inches(0.12), Inches(12.2), Inches(0.55),
        [[(title, 30, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(s, Inches(0.57), Inches(0.72), Inches(12.2), Inches(0.35),
            [[(kicker, 13, RGBColor(0xB9, 0xC6, 0xE6), False)]])


# ============================================================ Slide 1: Title
s = slide()
bg(s, DARK)
box(s, 0, Inches(2.55), SW, Inches(0.06), fill=BLUE)
box(s, Inches(4.44), Inches(2.55), Inches(4.45), Inches(0.06), fill=GREEN)
box(s, Inches(8.89), Inches(2.55), Inches(4.45), Inches(0.06), fill=RED)
txt(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.6),
    [[("State Machine & Per-State Control", 44, WHITE, True)]])
txt(s, Inches(0.82), Inches(4.15), Inches(11.7), Inches(0.8),
    [[("GreenPinkCamera Injectable Teardown", 24, RGBColor(0xB9, 0xC6, 0xE6), False)]])
txt(s, Inches(0.82), Inches(5.5), Inches(11.7), Inches(0.6),
    [[("Flexiv Rizon4 · 6-state FSM · primitive / Cartesian / joint-impedance control",
       15, GREY, False)]])

# ====================================== Slide 2: Two-layer controller framing
s = slide()
bg(s)
header(s, "The Controller Has Two Layers",
       "Describe every state with BOTH layers")

# Layer 1 card
box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), fill=LIGHT)
txt(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.5),
    [[("Layer 1 — Flexiv RDK Control Mode", 18, DARK, True)]])
l1 = [
    ("NRT_PRIMITIVE_EXECUTION", "Runs Flexiv primitives: MovePTP, MoveL, MoveC, "
     "InsertComp, FloatingCartesian, Contact, GraspComp, ZeroFTSensor", ACCENT),
    ("NRT_CARTESIAN_MOTION_FORCE", "Direct operational-space motion / force "
     "(grasp descent)", GREEN),
    ("NRT_JOINT_IMPEDANCE", "Joint hold: SendJointPosition + SetJointImpedance "
     "(park arm rigidly during gripper-only actions)", BLUE),
]
y = 2.2
for name, desc, c in l1:
    box(s, Inches(0.7), Inches(y), Inches(5.6), Inches(1.35), fill=WHITE, line=c, line_w=1.5)
    txt(s, Inches(0.85), Inches(y + 0.1), Inches(5.3), Inches(0.4),
        [[(name, 13.5, c, True)]])
    txt(s, Inches(0.85), Inches(y + 0.5), Inches(5.3), Inches(0.8),
        [[(desc, 11.5, GREY, False)]])
    y += 1.5

# Layer 2 card
box(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(5.6), fill=LIGHT)
txt(s, Inches(7.05), Inches(1.55), Inches(5.6), Inches(0.5),
    [[("Layer 2 — Control Paradigm", 18, DARK, True)]])
txt(s, Inches(7.05), Inches(1.95), Inches(5.6), Inches(0.35),
    [[("the Force / Joint / Operational language", 12, GREY, False)]])
l2 = [
    ("JOINT SPACE", BLUE, "MovePTP (planned in joint space, % vel scale), "
     "joint-impedance hold, the cap twist"),
    ("OPERATIONAL / CARTESIAN", GREEN, "MoveL, MoveC — commanded in m/s"),
    ("FORCE / COMPLIANCE", RED, "InsertComp, FloatingCartesian, Contact, "
     "GraspComp, Cartesian impedance, gripper force-target 80 N, vise close-to-force 5 kg"),
]
y = 2.45
for name, c, desc in l2:
    box(s, Inches(7.05), Inches(y), Inches(5.6), Inches(1.35), fill=c)
    txt(s, Inches(7.2), Inches(y + 0.12), Inches(5.3), Inches(0.4),
        [[(name, 14, WHITE, True)]])
    txt(s, Inches(7.2), Inches(y + 0.55), Inches(5.3), Inches(0.75),
        [[(desc, 11.5, WHITE, False)]])
    y += 1.5

# ============================================ Slide 3: State machine (anchor)
s = slide()
bg(s)
header(s, "The State Machine", "6 states · transition guards make it a real FSM")

states = ["Initialization", "Module\nLocalization", "Pick\nInjectable",
          "Load\nInjectable", "Ultrasonic\nCutting", "Component\nDisassembly"]
# color per dominant paradigm
scols = [BLUE, BLUE, RED, RED, ACCENT, GREEN]
n = len(states)
gap = Inches(0.16)
total_w = SW - Inches(1.0)
bw = Emu(int((total_w - gap * (n - 1)) / n))
x = Inches(0.5)
y = Inches(1.7)
bh = Inches(1.3)
for i, (st, c) in enumerate(zip(states, scols)):
    box(s, x, y, bw, bh, fill=c)
    txt(s, x, y, bw, bh, [[(st, 14, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        ax = Emu(int(x) + int(bw) + int(gap) // 2 - Inches(0.09))
        txt(s, ax, y, Inches(0.2), bh, [[("→", 22, GREY, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x = Emu(int(x) + int(bw) + int(gap))

# guards
box(s, Inches(0.5), Inches(3.4), Inches(12.33), Inches(3.55), fill=LIGHT)
txt(s, Inches(0.7), Inches(3.5), Inches(12), Inches(0.4),
    [[("Transition guards (from recipe safety gates)", 16, DARK, True)]])
guards = [
    ("→ enter vise area", "Arduino busy=false, faulted=false, blade_on=false, homed"),
    ("→ cut", "vise CLOSED, robot clear, rotary ≈ 0"),
    ("→ after cut", "CUT_HEIGHT=DONE, blade_on=false, x=z=0"),
    ("fault edge (any state)", "→ STOP_ALL (fault handler)"),
]
gy = 4.05
for cond, detail in guards:
    txt(s, Inches(0.8), Inches(gy), Inches(3.2), Inches(0.5),
        [[(cond, 14, DARK, True)]])
    txt(s, Inches(4.1), Inches(gy), Inches(8.5), Inches(0.5),
        [[(detail, 13.5, GREY, False)]])
    gy += 0.68

# ============================================ Slide 4: Controller legend
s = slide()
bg(s)
header(s, "Controller Legend", "Color code reused on every state slide")
legend = [
    (BLUE, "JOINT", "NRT_JOINT_IMPEDANCE", "MovePTP · joint hold · cap twist"),
    (GREEN, "OPERATIONAL", "NRT_CARTESIAN_MOTION_FORCE", "MoveL · MoveC (m/s)"),
    (RED, "FORCE / COMPLIANCE", "force-aware primitives",
     "InsertComp · FloatingCartesian · Contact · GraspComp · gripper 80 N · vise 5 kg"),
]
y = 1.55
for c, name, mode, prims in legend:
    box(s, Inches(0.5), Inches(y), Inches(0.5), Inches(1.5), fill=c)
    box(s, Inches(1.1), Inches(y), Inches(11.7), Inches(1.5), fill=LIGHT)
    txt(s, Inches(1.35), Inches(y + 0.15), Inches(4.5), Inches(0.5),
        [[(name, 20, c, True)]])
    txt(s, Inches(1.35), Inches(y + 0.78), Inches(5.0), Inches(0.5),
        [[(mode, 12.5, GREY, False)]])
    txt(s, Inches(6.0), Inches(y + 0.15), Inches(6.6), Inches(1.2),
        [[("Primitives:", 12.5, DARK, True)], [(prims, 14, DARK, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
    y += 1.72

# ============================================ Slides 5-10: per-state
# each: (title, kicker, sequence list[(label, color)], paradigm tags, mode, params)
state_data = [
    ("Initialization", "Bring robot, gripper and machine to a known safe state",
     [("gripper open", RED), ("MovePTP → Inter", BLUE), ("MovePTP → Home", BLUE)],
     [("Joint-space", BLUE), ("Gripper", RED)],
     "NRT_PRIMITIVE_EXECUTION",
     "MOVE_JNT_VEL_SCALE = 80%  ·  gripper 0.05 m/s"),

    ("Module Localization",
     "Re-localize the vise via ChArUco before any vise interaction",
     [("MovePTP → Vise-cali", BLUE), ("ChArUco detect (cali tag 1/2/3)", ACCENT),
      ("reload Vise.json", ACCENT), ("MovePTP → Home", BLUE)],
     [("Joint-space", BLUE), ("Vision (no closed-loop force)", ACCENT)],
     "NRT_PRIMITIVE_EXECUTION",
     "hand-eye camera_tcp.yaml  ·  tag_01_to_vise_tcp.json"),

    ("Pick Injectable",
     "Camera-aligned approach + contact-guided adaptive grasp",
     [("MovePTP → Plate", BLUE), ("Align Injectable (MoveL)", GREEN),
      ("ZeroFTSensor → Contact → GraspComp", RED), ("MoveL lift +20 cm", GREEN)],
     [("Operational", GREEN), ("Force (contact grasp)", RED)],
     "NRT_PRIMITIVE_EXECUTION  +  NRT_CARTESIAN_MOTION_FORCE",
     "align 0.02 m/s  ·  contact 0.05 m/s  ·  grip 80 N"),

    ("Load Injectable",
     "SHOWCASE STATE — cycles through all three RDK modes",
     [("MovePTP → above_vise", BLUE), ("InsertComp (−Z → TCP X)", RED),
      ("FloatingCartesian y/z/rx while vise → 5 kg", RED),
      ("joint hold", BLUE), ("gripper release", RED), ("MoveL retreat", GREEN)],
     [("Force/compliance", RED), ("→ Joint hold", BLUE), ("→ Operational", GREEN)],
     "ALL THREE: PRIMITIVE_EXECUTION · CARTESIAN_MOTION_FORCE · JOINT_IMPEDANCE",
     "INSERTCOMP_INSERT_VEL = 0.02  ·  maxContactForce  ·  float maxVel 0.2 m/s"),

    ("Ultrasonic Cutting",
     "Machine cuts while robot waits clear; then grip + twist off cap",
     [("Arduino CUT_HEIGHT (robot idle/clear)", ACCENT), ("MoveL down", GREEN),
      ("FloatingCartesian y/z/rx", RED), ("gripper 80 N", RED),
      ("twist about TCP X  0→+7→−10°", BLUE), ("MoveL lift", GREEN)],
     [("Machine", ACCENT), ("Force-float", RED), ("Joint-space twist", BLUE)],
     "NRT_PRIMITIVE_EXECUTION",
     "CUT z=134.2, x=110.5, deg=360  ·  twist scale 20%"),

    ("Component Disassembly",
     "Fast-mode pick-and-place of spring, plastic, shell & glass",
     [("MovePTP transit (×N)", BLUE), ("MoveL descend / lift 0.5 m/s", GREEN),
      ("gripper 80 N", RED), ("Arduino OPEN_VISE", ACCENT),
      ("MoveC dump arc", GREEN)],
     [("Joint transit", BLUE), ("Operational", GREEN), ("Force-grip", RED)],
     "NRT_PRIMITIVE_EXECUTION",
     "FRAME5_CARTESIAN_VEL = 0.5 m/s  ·  dump 0.03 m/s"),
]

for idx, (title, kicker, seq, tags, mode, params) in enumerate(state_data, 1):
    s = slide()
    bg(s)
    header(s, f"State {idx} — {title}", kicker)

    # paradigm tag chips
    cx = Inches(0.55)
    txt(s, Inches(0.55), Inches(1.25), Inches(2.2), Inches(0.35),
        [[("CONTROL:", 12, GREY, True)]])
    cx = Inches(1.85)
    for tname, tc in tags:
        w = Inches(0.42 + 0.105 * len(tname))
        box(s, cx, Inches(1.22), w, Inches(0.4), fill=tc)
        txt(s, cx, Inches(1.22), w, Inches(0.4), [[(tname, 12, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx = Emu(int(cx) + int(w) + Inches(0.12))

    # sequence flow
    txt(s, Inches(0.55), Inches(1.9), Inches(6), Inches(0.35),
        [[("Primitive / action sequence", 15, DARK, True)]])
    y = 2.4
    for i, (label, c) in enumerate(seq):
        box(s, Inches(0.7), Inches(y), Inches(7.4), Inches(0.62), fill=WHITE,
            line=c, line_w=2.0)
        box(s, Inches(0.7), Inches(y), Inches(0.12), Inches(0.62), fill=c)
        txt(s, Inches(1.0), Inches(y), Inches(7.0), Inches(0.62),
            [[(f"{i+1}.  ", 13, c, True), (label, 13.5, DARK, False)]],
            anchor=MSO_ANCHOR.MIDDLE)
        y += 0.72

    # right panel: mode + params
    box(s, Inches(8.45), Inches(2.4), Inches(4.4), Inches(2.1), fill=LIGHT)
    txt(s, Inches(8.65), Inches(2.55), Inches(4.0), Inches(0.4),
        [[("RDK control mode", 13, GREY, True)]])
    txt(s, Inches(8.65), Inches(2.95), Inches(4.0), Inches(1.4),
        [[(mode, 13.5, DARK, True)]])

    box(s, Inches(8.45), Inches(4.65), Inches(4.4), Inches(2.0), fill=LIGHT)
    txt(s, Inches(8.65), Inches(4.8), Inches(4.0), Inches(0.4),
        [[("Defining parameters", 13, GREY, True)]])
    txt(s, Inches(8.65), Inches(5.2), Inches(4.0), Inches(1.4),
        [[(params, 13.5, DARK, False)]])

# ============================================ Slide 11: takeaways
s = slide()
bg(s, DARK)
txt(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8),
    [[("Two Things To Remember", 34, WHITE, True)]])
box(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.2), fill=RGBColor(0x2A, 0x37, 0x57))
txt(s, Inches(1.1), Inches(1.95), Inches(11.1), Inches(1.8),
    [[("1 · Color-code by control paradigm", 22, RGBColor(0x9D, 0xC4, 0xFF), True)],
     [("Blue = joint, green = operational, red = force — consistent across the "
       "FSM and every state slide. Pick & Load read as force-rich; Disassembly "
       "is mostly kinematic.", 16, RGBColor(0xD9, 0xE2, 0xF2), False)]])
box(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.4), fill=RGBColor(0x2A, 0x37, 0x57))
txt(s, Inches(1.1), Inches(4.35), Inches(11.1), Inches(2.0),
    [[("2 · Load Injectable is the showcase state", 22, RGBColor(0xFF, 0xB3, 0xB3), True)],
     [("It is the only state that cycles through all three RDK modes: compliant "
       "InsertComp → floating force → joint-impedance hold → Cartesian retreat. "
       "This is your 'we understand control architecture' slide.", 16,
       RGBColor(0xD9, 0xE2, 0xF2), False)]])

out = "/home/src0/flexiv_rdk/project/state_machine_controller_slides.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
